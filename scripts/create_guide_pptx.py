"""홈케어커넥트 이용자 설명서 파워포인트 생성"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 색상 정의
PRIMARY = RGBColor(0x00, 0x20, 0x45)
SECONDARY = RGBColor(0x00, 0x6A, 0x63)
SURFACE = RGBColor(0xF7, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x18, 0x1C, 0x1E)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_BG = RGBColor(0xF1, 0xF4, 0xF6)
TEAL_LIGHT = RGBColor(0xE0, 0xF7, 0xF5)
ERROR = RGBColor(0xBA, 0x1A, 0x1A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=SURFACE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def add_bullet_text(slide, left, top, width, height, items, font_size=16, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
        p.level = 0
    return txBox

# ═══════════════════════════════════════════
# 슬라이드 1: 표지
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, PRIMARY)

# 장식 원
shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-1), Inches(5), Inches(5))
shape.fill.solid()
shape.fill.fore_color.rgb = SECONDARY
shape.fill.fore_color.brightness = 0.3
shape.line.fill.background()

add_text(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1), '홈케어커넥트', 48, WHITE, True, PP_ALIGN.LEFT)
add_text(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(0.8), 'AI 기반 방문치료 매칭 & 운영 플랫폼', 24, RGBColor(0xA0, 0xF1, 0xE8), False, PP_ALIGN.LEFT)
add_text(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.5), '이용자 설명서', 20, RGBColor(0xB8, 0xC5, 0xD6), False, PP_ALIGN.LEFT)
add_text(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.5), '주식회사 루미브리즈  |  2026년 3월', 14, RGBColor(0x80, 0x90, 0xA0), False, PP_ALIGN.LEFT)

# ═══════════════════════════════════════════
# 슬라이드 2: 목차
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SURFACE)

add_text(slide, Inches(1.5), Inches(0.8), Inches(10), Inches(0.8), '목차', 36, PRIMARY, True)

toc_items = [
    '1.  홈케어커넥트란?',
    '2.  보호자(환자) — 이렇게 사용하세요',
    '3.  간호사 — 이렇게 사용하세요',
    '4.  병원/기관 관리자 — 이렇게 사용하세요',
    '5.  AI 돌봄 도우미',
    '6.  스마트 복약 관리',
    '7.  자주 묻는 질문',
]

for i, item in enumerate(toc_items):
    y = Inches(2.0 + i * 0.65)
    shape = add_shape(slide, Inches(1.5), y, Inches(10), Inches(0.55), WHITE if i % 2 == 0 else LIGHT_BG, 0.05)
    add_text(slide, Inches(1.8), y + Emu(Inches(0.08).emu), Inches(9), Inches(0.4), item, 18, PRIMARY, False)

# ═══════════════════════════════════════════
# 슬라이드 3: 홈케어커넥트란?
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SURFACE)

add_text(slide, Inches(1.5), Inches(0.8), Inches(10), Inches(0.8), '1. 홈케어커넥트란?', 36, PRIMARY, True)
add_text(slide, Inches(1.5), Inches(1.7), Inches(10), Inches(0.5), 'AI 기반 방문치료 매칭 & 운영 SaaS 플랫폼', 18, GRAY)

# 4개 핵심 기능 카드
features = [
    ('🔍', 'AI 기관 매칭', '지역·전문성 기반\n최적 기관 자동 추천'),
    ('📊', '건강 모니터링', '바이탈 추적 +\nAI 이상징후 감지'),
    ('💊', '스마트 복약', '자동 알림 +\nAI 복약지도'),
    ('🤖', 'AI 돌봄 도우미', '음성 대화로\n건강 관리'),
]

for i, (icon, title, desc) in enumerate(features):
    x = Inches(1.5 + i * 2.8)
    card = add_shape(slide, x, Inches(2.8), Inches(2.5), Inches(2.5), WHITE, 0.08)
    add_text(slide, x + Inches(0.3), Inches(3.0), Inches(2), Inches(0.5), icon, 36, DARK)
    add_text(slide, x + Inches(0.3), Inches(3.6), Inches(2), Inches(0.4), title, 18, PRIMARY, True)
    add_text(slide, x + Inches(0.3), Inches(4.1), Inches(2), Inches(0.8), desc, 14, GRAY)

# 하단 설명
add_text(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(1),
    '보호자, 간호사, 병원, 플랫폼 관리자 — 4가지 역할을 하나의 앱에서 모두 지원합니다.\n로그인하면 역할에 맞는 화면이 자동으로 표시됩니다.', 16, GRAY)

# ═══════════════════════════════════════════
# 슬라이드 4: 보호자 - 메인 기능
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SURFACE)

# 좌측 제목 영역
add_shape(slide, Inches(0), Inches(0), Inches(4.5), Inches(7.5), PRIMARY)
add_text(slide, Inches(0.8), Inches(1), Inches(3.5), Inches(0.5), '보호자/환자', 14, RGBColor(0xA0, 0xF1, 0xE8), True)
add_text(slide, Inches(0.8), Inches(1.5), Inches(3.5), Inches(1), '이렇게\n사용하세요', 36, WHITE, True)
add_text(slide, Inches(0.8), Inches(3.0), Inches(3.5), Inches(2),
    '어르신의 건강을\nAI가 24시간 모니터링\n하고 최적의 방문치료\n기관을 매칭해드립니다.', 16, RGBColor(0xB8, 0xC5, 0xD6))

# 우측 기능 리스트
funcs = [
    ('홈 대시보드', '오늘의 방문 일정, 건강 점수, 케어 타임라인을 한눈에 확인'),
    ('AI 기관 매칭', '지역·서비스·시간 선택 → AI가 최적 기관 3곳 추천'),
    ('방문 일정', '주간/월간 캘린더로 방문 스케줄 관리'),
    ('방문 기록', '바이탈 추이 차트 + 간호 노트 + 수행 내역 확인'),
    ('환자 등록/관리', '어르신 정보 등록, 요양등급, 진단명, 필요 서비스 설정'),
    ('AI 리포트', 'AI가 분석한 주간 건강 리포트 + 의사 소견 확인'),
    ('리뷰 작성', '이용한 기관에 별점 + 세부 평가 작성'),
    ('알림', '방문 알림, 매칭 결과, 레드플래그 경고 수신'),
]

for i, (title, desc) in enumerate(funcs):
    y = Inches(0.7 + i * 0.8)
    add_shape(slide, Inches(5), y, Inches(7.8), Inches(0.7), WHITE if i % 2 == 0 else LIGHT_BG, 0.03)
    add_text(slide, Inches(5.3), y + Emu(Inches(0.05).emu), Inches(2.2), Inches(0.3), f'▸ {title}', 15, SECONDARY, True)
    add_text(slide, Inches(5.3), y + Emu(Inches(0.35).emu), Inches(7.2), Inches(0.3), desc, 13, GRAY)

# ═══════════════════════════════════════════
# 슬라이드 5: 보호자 - 사용 흐름
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SURFACE)

add_text(slide, Inches(1.5), Inches(0.8), Inches(10), Inches(0.8), '보호자 — 처음 사용하기', 36, PRIMARY, True)

steps = [
    ('1', '회원가입', '앱 실행 → [시작하기] → 회원가입\n이름, 전화번호, 이메일 입력\n역할: "보호자" 선택'),
    ('2', '환자 등록', '마이페이지 → 환자 관리 → 등록\n어르신 이름, 생년월일, 성별\n요양등급, 주소, 필요 서비스 입력'),
    ('3', '기관 매칭', '매칭 탭 → [매칭 시작하기]\n서비스, 지역, 시간 선택\nAI가 3곳 추천 → 기관 선택'),
    ('4', '일정 확인', '일정 탭에서 방문 스케줄 확인\n방문 당일 푸시 알림 수신\n방문 완료 후 기록 확인'),
]

for i, (num, title, desc) in enumerate(steps):
    x = Inches(1.2 + i * 2.9)
    # 번호 원
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.8), Inches(2.2), Inches(0.6), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SECONDARY
    shape.line.fill.background()
    add_text(slide, x + Inches(0.8), Inches(2.25), Inches(0.6), Inches(0.5), num, 22, WHITE, True, PP_ALIGN.CENTER)

    # 화살표 (마지막 제외)
    if i < 3:
        add_text(slide, x + Inches(2.3), Inches(2.3), Inches(0.5), Inches(0.4), '→', 24, GRAY, False, PP_ALIGN.CENTER)

    # 카드
    add_shape(slide, x, Inches(3.2), Inches(2.6), Inches(3.2), WHITE, 0.06)
    add_text(slide, x + Inches(0.3), Inches(3.4), Inches(2), Inches(0.4), title, 20, PRIMARY, True)
    add_text(slide, x + Inches(0.3), Inches(3.9), Inches(2.2), Inches(2.2), desc, 14, GRAY)

# ═══════════════════════════════════════════
# 슬라이드 6: 간호사 - 메인 기능
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SURFACE)

add_shape(slide, Inches(0), Inches(0), Inches(4.5), Inches(7.5), SECONDARY)
add_text(slide, Inches(0.8), Inches(1), Inches(3.5), Inches(0.5), '간호사', 14, RGBColor(0xA0, 0xF1, 0xE8), True)
add_text(slide, Inches(0.8), Inches(1.5), Inches(3.5), Inches(1), '이렇게\n사용하세요', 36, WHITE, True)
add_text(slide, Inches(0.8), Inches(3.0), Inches(3.5), Inches(2),
    '오늘의 방문 일정부터\n바이탈 기록, 레드플래그\n알림까지 — 현장 업무를\n스마트하게 관리합니다.', 16, RGBColor(0xD0, 0xFF, 0xF8))

funcs = [
    ('오늘 일정', 'LIVE 방문 카드 + 남은 방문 리스트 + 실시간 환자 바이탈'),
    ('방문 수행 (5단계)', '① 체크인(GPS) → ② 바이탈 입력 → ③ 수행 체크 → ④ 메모/사진/음성 → ⑤ 체크아웃'),
    ('담당 환자', '환자별 바이탈 추이, 처방약, 컨디션 체크 기록 확인'),
    ('레드플래그', 'AI가 감지한 이상징후 알림 — 심각도별 색상 (빨강/주황/노랑)'),
    ('AI 어시스턴트', '"오늘 브리핑 해줘" → 전체 일정 + 주의환자 요약 음성 안내'),
    ('월간 통계', '방문 횟수, 완료율, 총 시간, 평균 시간 통계'),
    ('오프라인 모드', '인터넷 없는 곳에서도 기록 → 온라인 복귀 시 자동 동기화'),
]

for i, (title, desc) in enumerate(funcs):
    y = Inches(0.7 + i * 0.9)
    add_shape(slide, Inches(5), y, Inches(7.8), Inches(0.8), WHITE if i % 2 == 0 else LIGHT_BG, 0.03)
    add_text(slide, Inches(5.3), y + Emu(Inches(0.05).emu), Inches(2.5), Inches(0.3), f'▸ {title}', 15, SECONDARY, True)
    add_text(slide, Inches(5.3), y + Emu(Inches(0.38).emu), Inches(7.2), Inches(0.35), desc, 13, GRAY)

# ═══════════════════════════════════════════
# 슬라이드 7: 간호사 - 방문 수행 5단계
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SURFACE)

add_text(slide, Inches(1.5), Inches(0.8), Inches(10), Inches(0.8), '간호사 — 방문 수행 5단계', 36, PRIMARY, True)
add_text(slide, Inches(1.5), Inches(1.7), Inches(10), Inches(0.5), '환자 방문 시 아래 순서로 기록합니다. 각 단계는 앱에서 안내됩니다.', 16, GRAY)

visit_steps = [
    ('① 체크인', 'GPS 자동 인식\n도착 시간 기록\n위치 확인 후 [체크인]'),
    ('② 바이탈', '혈압, 심박, 체온\n혈당, 산소포화도\n정상/경고 자동 표시'),
    ('③ 수행 체크', '계획된 케어 항목\n체크리스트로 확인\n전체 선택 가능'),
    ('④ 메모', '간호 노트 작성\n사진 촬영, 음성 메모\n보호자 메시지 입력'),
    ('⑤ 체크아웃', '방문 요약 확인\n소요 시간 자동 계산\n[완료] → AI 분석 시작'),
]

for i, (title, desc) in enumerate(visit_steps):
    x = Inches(0.8 + i * 2.5)
    # 상단 연결선
    if i < 4:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(2.2), Inches(3.0), Inches(0.5), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = SECONDARY
        line.line.fill.background()
    # 카드
    card_color = SECONDARY if i == 0 or i == 4 else WHITE
    text_color = WHITE if card_color == SECONDARY else PRIMARY
    add_shape(slide, x, Inches(2.5), Inches(2.2), Inches(3.5), card_color, 0.06)
    add_text(slide, x + Inches(0.2), Inches(2.7), Inches(1.8), Inches(0.5), title, 20, text_color, True)
    desc_color = RGBColor(0xD0, 0xFF, 0xF8) if card_color == SECONDARY else GRAY
    add_text(slide, x + Inches(0.2), Inches(3.4), Inches(1.8), Inches(2.3), desc, 14, desc_color)

# ═══════════════════════════════════════════
# 슬라이드 8: 병원/기관 관리자
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SURFACE)

add_shape(slide, Inches(0), Inches(0), Inches(4.5), Inches(7.5), RGBColor(0x32, 0x1B, 0x00))
add_text(slide, Inches(0.8), Inches(1), Inches(3.5), Inches(0.5), '병원/기관 관리자', 14, RGBColor(0xFF, 0xE0, 0xB2), True)
add_text(slide, Inches(0.8), Inches(1.5), Inches(3.5), Inches(1), '이렇게\n사용하세요', 36, WHITE, True)
add_text(slide, Inches(0.8), Inches(3.0), Inches(3.5), Inches(2),
    '환자 관리, 직원 배정,\n일정, 통계, 수납,\n건보 청구까지 —\n기관 운영을 한 곳에서.', 16, RGBColor(0xFF, 0xE0, 0xB2))

funcs = [
    ('대시보드', '오늘의 KPI (대기/진행/완료), 레드플래그, 최근 요청, 방문 일정 요약'),
    ('환자 관리', '소속 환자 검색/필터, 상세 정보, 바이탈 이력, 서비스 플랜'),
    ('직원 관리', '간호사/치료사 목록, 직원 초대(이메일), 담당 환자 배정'),
    ('일정 관리', '전체 직원 방문 일정 캘린더, 날짜별 상세'),
    ('서비스 요청', '보호자 매칭 요청 수신 → [수락]/[거절] → 간호사 배정'),
    ('의사 방문', '소견 입력 → AI가 어르신 눈높이로 자동 변환'),
    ('수납/건보청구', '수납 현황 관리, 건강보험 청구 자료 조회'),
    ('통계', '월별 방문수, 완료율, 총 시간, 매출 현황'),
]

for i, (title, desc) in enumerate(funcs):
    y = Inches(0.5 + i * 0.85)
    add_shape(slide, Inches(5), y, Inches(7.8), Inches(0.75), WHITE if i % 2 == 0 else LIGHT_BG, 0.03)
    add_text(slide, Inches(5.3), y + Emu(Inches(0.05).emu), Inches(2.2), Inches(0.3), f'▸ {title}', 15, RGBColor(0x32, 0x1B, 0x00), True)
    add_text(slide, Inches(5.3), y + Emu(Inches(0.38).emu), Inches(7.2), Inches(0.3), desc, 13, GRAY)

# ═══════════════════════════════════════════
# 슬라이드 9: AI 돌봄 도우미
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SURFACE)

add_text(slide, Inches(1.5), Inches(0.8), Inches(10), Inches(0.8), '5. AI 돌봄 도우미', 36, PRIMARY, True)
add_text(slide, Inches(1.5), Inches(1.7), Inches(10), Inches(0.5), '어르신과 간호사를 위한 AI 대화 비서', 18, GRAY)

# 환자용
add_shape(slide, Inches(1), Inches(2.5), Inches(5.5), Inches(4.5), WHITE, 0.06)
add_text(slide, Inches(1.3), Inches(2.7), Inches(5), Inches(0.5), '💬 보호자/환자 도우미', 20, PRIMARY, True)
add_text(slide, Inches(1.3), Inches(3.3), Inches(5), Inches(0.3), '다정한 동네 간호사 페르소나', 14, SECONDARY, True)

patient_features = [
    '• "오늘 일정 알려줘" → 방문 스케줄 음성 안내',
    '• "약 먹을 시간이야?" → 복약 스케줄 확인',
    '• "몸이 안 좋아요" → 컨디션 체크 + 필요 시 간호사 알림',
    '• 식사 사진 → AI 영양 분석 (단백질 섭취 강조)',
    '• 마이크 버튼으로 음성 대화 (어르신 친화)',
    '• 큰 글씨(18px) + TTS 읽어주기 기능',
]
add_bullet_text(slide, Inches(1.3), Inches(3.8), Inches(5), Inches(3), patient_features, 14, GRAY)

# 간호사용
add_shape(slide, Inches(7), Inches(2.5), Inches(5.5), Inches(4.5), WHITE, 0.06)
add_text(slide, Inches(7.3), Inches(2.7), Inches(5), Inches(0.5), '🩺 간호사 어시스턴트', 20, PRIMARY, True)
add_text(slide, Inches(7.3), Inches(3.3), Inches(5), Inches(0.3), '간결한 업무 브리핑 비서', 14, SECONDARY, True)

nurse_features = [
    '• "오늘 브리핑 해줘" → 전체 일정 + 주의환자 요약',
    '• "다음 환자 알려줘" → 이동 중 음성 브리핑',
    '• "주의 환자 있어?" → 레드플래그 + 복약 미이행 확인',
    '• 환자별 바이탈 트렌드 + 컨디션 체크 결과 통합',
    '• 처방약 변경 사항 하이라이트',
    '• 임상 가이드라인 RAG 검색',
]
add_bullet_text(slide, Inches(7.3), Inches(3.8), Inches(5), Inches(3), nurse_features, 14, GRAY)

# ═══════════════════════════════════════════
# 슬라이드 10: 스마트 복약 관리
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SURFACE)

add_text(slide, Inches(1.5), Inches(0.8), Inches(10), Inches(0.8), '6. 스마트 복약 관리', 36, PRIMARY, True)

# 흐름도
flow_steps = [
    ('의사가 처방 입력', '병원 앱에서\n약품명, 용량, 횟수 입력'),
    ('AI 복약지도 생성', 'e약은요 API 조회 +\nDUR 병용금기 체크 +\n어르신 눈높이 변환'),
    ('자동 알람 발송', '복용 시간 → 푸시 알림\n"어머님, 혈압약\n드실 시간이에요~"'),
    ('복용 확인/미이행', '환자가 [복용완료] 탭\n30분 미응답 → 재알림\n60분 미응답 → 간호사 알림'),
]

for i, (title, desc) in enumerate(flow_steps):
    x = Inches(1 + i * 3.1)
    if i < 3:
        add_text(slide, x + Inches(2.5), Inches(3.5), Inches(0.5), Inches(0.5), '→', 28, SECONDARY, True, PP_ALIGN.CENTER)
    card_bg = SECONDARY if i == 2 else WHITE
    txt_color = WHITE if i == 2 else PRIMARY
    desc_color = RGBColor(0xD0, 0xFF, 0xF8) if i == 2 else GRAY
    add_shape(slide, x, Inches(2.5), Inches(2.8), Inches(3), card_bg, 0.06)
    add_text(slide, x + Inches(0.3), Inches(2.7), Inches(2.2), Inches(0.5), title, 18, txt_color, True)
    add_text(slide, x + Inches(0.3), Inches(3.3), Inches(2.2), Inches(2), desc, 14, desc_color)

add_text(slide, Inches(1.5), Inches(6), Inches(10), Inches(0.8),
    '⚠️ 병용금기 발견 시 → 의사/간호사에게 자동 알림  |  모든 복약지도는 공공 API(e약은요, DUR) 근거 기반', 14, ERROR)

# ═══════════════════════════════════════════
# 슬라이드 11: 자주 묻는 질문
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SURFACE)

add_text(slide, Inches(1.5), Inches(0.8), Inches(10), Inches(0.8), '7. 자주 묻는 질문 (FAQ)', 36, PRIMARY, True)

faqs = [
    ('Q. 앱은 어디서 다운받나요?', 'A. App Store(iPhone) 또는 Google Play(Android)에서 "홈케어커넥트"를 검색하세요.'),
    ('Q. 하나의 앱으로 모든 역할을 사용할 수 있나요?', 'A. 네, 로그인 시 등록된 역할에 따라 자동으로 맞는 화면이 표시됩니다.'),
    ('Q. AI 도우미가 의료 진단을 하나요?', 'A. 아닙니다. AI는 건강 정보 안내와 이상징후 감지만 수행하며, 의료 판단은 반드시 담당 의료진을 따르세요.'),
    ('Q. 인터넷이 안 되는 곳에서도 사용 가능한가요?', 'A. 간호사 앱은 오프라인에서도 방문 기록이 가능합니다. 인터넷 연결 시 자동 동기화됩니다.'),
    ('Q. 개인정보는 안전한가요?', 'A. 모든 데이터는 암호화 저장되며, RLS(행 수준 보안)로 본인 + 담당자만 접근 가능합니다.'),
    ('Q. 복약 알림은 어떻게 오나요?', 'A. 의사가 처방을 입력하면 복용 시간에 맞춰 자동으로 푸시 알림이 발송됩니다.'),
]

for i, (q, a) in enumerate(faqs):
    y = Inches(1.8 + i * 0.9)
    add_shape(slide, Inches(1.5), y, Inches(10.3), Inches(0.8), WHITE if i % 2 == 0 else LIGHT_BG, 0.03)
    add_text(slide, Inches(1.8), y + Emu(Inches(0.05).emu), Inches(9.5), Inches(0.3), q, 15, PRIMARY, True)
    add_text(slide, Inches(1.8), y + Emu(Inches(0.4).emu), Inches(9.5), Inches(0.3), a, 13, GRAY)

# ═══════════════════════════════════════════
# 슬라이드 12: 마무리
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)

shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(4), Inches(6), Inches(6))
shape.fill.solid()
shape.fill.fore_color.rgb = SECONDARY
shape.fill.fore_color.brightness = 0.3
shape.line.fill.background()

add_text(slide, Inches(1.5), Inches(2), Inches(10), Inches(1),
    '기술로 잇는 따뜻한 돌봄', 44, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.5),
    '홈케어커넥트와 함께 더 나은 방문치료를 경험하세요', 20, RGBColor(0xA0, 0xF1, 0xE8), False, PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(5), Inches(10), Inches(0.5),
    '문의: support@lumibreeze.kr  |  www.homecareconnect.kr', 14, RGBColor(0x80, 0x90, 0xA0), False, PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(5.6), Inches(10), Inches(0.5),
    '주식회사 루미브리즈', 16, RGBColor(0xB8, 0xC5, 0xD6), False, PP_ALIGN.CENTER)

# 저장
output_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), '홈케어커넥트_이용자설명서.pptx')
prs.save(output_path)
print(f'저장 완료: {output_path}')
